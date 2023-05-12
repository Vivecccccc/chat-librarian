import os
import shutil
from datetime import datetime
from io import BufferedReader
from typing import Any, Dict, Optional, Tuple
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import docx
import csv
import pptx
import openpyxl
from handler.utils import hash_string

from models.document import DocumentMetadata, DocumentVersion, SingleDocument
from models.generic import Source

async def get_document_from_file(file: UploadFile) -> SingleDocument:
    extracted_text, meta = await extract_data_from_form_file(file)
    content_hash = hash_string(extracted_text)
    version = DocumentVersion(version_id=content_hash,
                              version_url=meta["filepath"],
                              modified_at=meta["modified_at"])
    # get metadata
    metadata = DocumentMetadata(source=Source.document,
                                created_at=meta["created_at"],
                                created_by=meta["created_by"],
                                version=version)
    metadata_hash = hash(metadata)
    doc_id = hash_string(str(metadata_hash))
    doc = SingleDocument(doc_id=doc_id, 
                         text=extracted_text, 
                         metadata=metadata)

    return doc

def extract_data_from_filepath(filepath: str, mimetype: Optional[str] = None) -> Tuple[str, Dict[str, Optional[Any]]]:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    # Open the file in binary mode
    file = open(filepath, "rb")
    extracted_text, created_at, created_by, modified_at = extract_data_from_file(file, mimetype)

    return extracted_text, {"created_at": created_at, "created_by": created_by, "modified_at": modified_at}

def extract_data_from_file(file: BufferedReader, mimetype: str) -> Tuple[str, Optional[datetime], Optional[str], Optional[datetime]]:
    created_at, created_by, modified_at = None, None, None
    if mimetype == "application/pdf":
        extracted_text, created_at, created_by, modified_at = _textify_pdf(file)
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        extracted_text, created_at, created_by, modified_at = _textify_docx(file)
    elif mimetype == "text/csv":
        extracted_text = _textify_csv(file)
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        extracted_text, created_at, created_by, modified_at = _textify_pptx(file)
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        extracted_text, created_at, created_by, modified_at = _textify_xlsx(file)
    else:
        # Unsupported file type
        file.close()
        raise ValueError("Unsupported file type: {}".format(mimetype))

    file.close()
    return extracted_text, created_at, created_by, modified_at

def _textify_xlsx(file: BufferedReader):
    wb = openpyxl.load_workbook(file)
    created_by = wb.properties.creator
    created_at = wb.properties.created
    modified_at = wb.properties.modified
    extracted_text = ""
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell is not None)
            if row_text:
                extracted_text += row_text + "\n"
    return extracted_text, created_at, created_by, modified_at

def _textify_pdf(file: BufferedReader):
    reader = PdfReader(file)
    created_by = reader.metadata.author
    created_at = reader.metadata.creation_date
    modified_at = reader.metadata.modification_date
    extracted_text = ""
    for page in reader.pages:
        extracted_text += page.extract_text()
    return extracted_text, created_at, created_by, modified_at

def _textify_docx(file: BufferedReader):
    doc = docx.Document(file)
    created_by = doc.core_properties.author
    created_at = doc.core_properties.created
    modified_at = doc.core_properties.modified
    extracted_text = docx2txt.process(file)
    return extracted_text, created_at, created_by, modified_at

def _textify_csv(file: BufferedReader):
    extracted_text = ""
    decoded_buffer = (line.decode("utf-8") for line in file)
    reader = csv.reader(decoded_buffer)
    for row in reader:
        extracted_text += " ".join(row) + "\n"
    return extracted_text

def _textify_pptx(file: BufferedReader):
    extracted_text = ""
    presentation = pptx.Presentation(file)
    created_by = presentation.core_properties.author
    created_at = presentation.core_properties.created
    modified_at = presentation.core_properties.modified
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        extracted_text += run.text + " "
                extracted_text += "\n"
    return extracted_text, created_at, created_by, modified_at

# Extract text from a file based on its mimetype
async def extract_data_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    print(f"mimetype: {mimetype}")
    print(f"file.file: {file.file}")
    print("file: ", file)

    filename = file.filename

    temp_dir = "tmp/.uploaded"
    os.makedirs(temp_dir, exist_ok=True)
    temp_file_path = os.path.join(temp_dir, filename)
    meta = {"filepath": filename, "created_at": None, "created_by": None, "modified_at": None}
    # write the file to a temporary locatoin
    with open(temp_file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    try:
        extracted_text, _meta = extract_data_from_filepath(temp_file_path, mimetype)
        meta.update(_meta)
    except Exception as e:
        print(f"Error: {e}")
        os.remove(temp_file_path)
        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text, meta